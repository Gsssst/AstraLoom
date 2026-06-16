from io import BytesIO
from types import SimpleNamespace
from uuid import UUID
import base64

import pytest
from pydantic import BaseModel

from app.services import chat_agent_tools
from app.services.chat_agent_tools import (
    AddToFolderArgs,
    ChatAgentRuntimeState,
    ChatAgentToolRuntime,
    ChatToolCall,
    ChatToolDefinition,
    ChatToolObservation,
    ChatToolRegistry,
    CreateResearchProjectArgs,
    ImportPaperArgs,
    chat_tool_confirmation_token,
    default_chat_tool_registry,
    deterministic_chat_tool_plan,
)
from app.db.models.paper import PaperFolderItem, UserPaper
from app.db.models.research import ResearchProject
from app.services.research_skills import ResearchSkillRunResult, get_research_skill


class EchoArgs(BaseModel):
    text: str


class _ScalarResult:
    def __init__(self, rows):
        self.rows = rows

    def scalars(self):
        return self

    def all(self):
        return list(self.rows)

    def scalar_one_or_none(self):
        return self.rows[0] if self.rows else None


class _FakeToolDb:
    def __init__(self, responses):
        self.responses = list(responses)
        self.added = []
        self.commit_count = 0
        self.refresh_count = 0

    async def execute(self, _query):
        return _ScalarResult(self.responses.pop(0) if self.responses else [])

    def add(self, item):
        self.added.append(item)

    async def commit(self):
        self.commit_count += 1

    async def refresh(self, item):
        self.refresh_count += 1
        if getattr(item, "id", None) is None:
            item.id = "project-1"


def test_registry_schema_export_includes_side_effect_policy():
    registry = default_chat_tool_registry()
    schemas = {schema["name"]: schema for schema in registry.schemas()}

    assert {
        "search_papers",
        "search_library",
        "import_paper",
        "read_pdf",
        "extract_docx",
        "extract_pptx",
        "run_skill",
        "add_to_folder",
        "create_research_project",
    } <= set(schemas)
    assert schemas["import_paper"]["side_effect"] is True
    assert schemas["add_to_folder"]["side_effect"] is True
    assert schemas["create_research_project"]["side_effect"] is True
    assert schemas["search_papers"]["side_effect"] is False
    assert schemas["read_pdf"]["side_effect"] is False
    assert schemas["extract_docx"]["side_effect"] is False
    assert schemas["extract_pptx"]["side_effect"] is False
    assert schemas["run_skill"]["side_effect"] is False


@pytest.mark.asyncio
async def test_unknown_tool_is_rejected_with_allowed_tools():
    registry = default_chat_tool_registry()
    state = ChatAgentRuntimeState(user_query="test")

    observation = await registry.execute(ChatToolCall(tool="missing_tool"), state)

    assert observation.status == "rejected"
    assert "search_papers" in observation.details["allowed_tools"]


@pytest.mark.asyncio
async def test_invalid_arguments_do_not_execute():
    called = False

    async def _executor(args, state):
        nonlocal called
        called = True
        return ChatToolObservation(tool="echo", status="completed")

    registry = ChatToolRegistry()
    registry.register(ChatToolDefinition(name="echo", label="Echo", args_model=EchoArgs, executor=_executor))

    observation = await registry.execute(ChatToolCall(tool="echo", arguments={}), ChatAgentRuntimeState(user_query="x"))

    assert observation.status == "rejected"
    assert called is False


@pytest.mark.asyncio
async def test_runtime_stops_at_max_steps_and_records_trace():
    async def _executor(args, state):
        return ChatToolObservation(tool="echo", status="completed", summary=args.text, result_count=1)

    registry = ChatToolRegistry()
    registry.register(ChatToolDefinition(name="echo", label="Echo", args_model=EchoArgs, executor=_executor))
    runtime = ChatAgentToolRuntime(registry, max_steps=1)
    state = await runtime.run(
        ChatAgentRuntimeState(user_query="x"),
        [
            ChatToolCall(tool="echo", arguments={"text": "one"}),
            ChatToolCall(tool="echo", arguments={"text": "two"}),
        ],
    )

    assert state.stop_reason == "max_steps"
    assert len(state.observations) == 1
    assert [event.status for event in state.trace_events] == ["running", "completed"]


@pytest.mark.asyncio
async def test_side_effect_tool_waits_for_exact_confirmation():
    registry = default_chat_tool_registry()
    args = {
        "source": "arxiv",
        "remote_id": "2601.00001",
        "auto_download": False,
        "title": "arXiv:2601.00001",
    }
    normalized_args = ImportPaperArgs.model_validate(args).model_dump()

    observation = await registry.execute(
        ChatToolCall(tool="import_paper", arguments=args),
        ChatAgentRuntimeState(user_query="import"),
        allow_side_effects=False,
    )

    assert observation.status == "waiting_confirmation"
    assert observation.details["confirmation_token"] == chat_tool_confirmation_token("import_paper", normalized_args)


@pytest.mark.asyncio
async def test_library_side_effect_tools_wait_for_confirmation():
    registry = default_chat_tool_registry()
    folder_args = AddToFolderArgs(
        folder_id="11111111-1111-1111-1111-111111111111",
        paper_ids=["22222222-2222-2222-2222-222222222222"],
    ).model_dump()
    project_args = CreateResearchProjectArgs(
        name="Video Grounding",
        description="seed project",
        paper_ids=["22222222-2222-2222-2222-222222222222"],
    ).model_dump()

    folder_observation = await registry.execute(
        ChatToolCall(tool="add_to_folder", arguments=folder_args),
        ChatAgentRuntimeState(user_query="add to folder"),
        allow_side_effects=False,
    )
    project_observation = await registry.execute(
        ChatToolCall(tool="create_research_project", arguments=project_args),
        ChatAgentRuntimeState(user_query="create project"),
        allow_side_effects=False,
    )

    assert folder_observation.status == "waiting_confirmation"
    assert folder_observation.details["confirmation_token"] == chat_tool_confirmation_token("add_to_folder", folder_args)
    assert project_observation.status == "waiting_confirmation"
    assert project_observation.details["confirmation_token"] == chat_tool_confirmation_token("create_research_project", project_args)


@pytest.mark.asyncio
async def test_search_papers_returns_references_and_context(monkeypatch):
    async def _fake_search(query, *, source, max_results, **kwargs):
        return [
            chat_agent_tools.PaperResult(
                title="Video Grounding Paper",
                authors=["A"],
                abstract="A useful paper.",
                year=2026,
                arxiv_id="2601.00001",
                source="arxiv",
                source_url="https://arxiv.org/abs/2601.00001",
                pdf_url="https://arxiv.org/pdf/2601.00001",
                metadata={"remote_id": "2601.00001"},
            )
        ]

    monkeypatch.setattr(chat_agent_tools, "search_scholarly_papers", _fake_search)

    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        ChatAgentRuntimeState(user_query="find papers"),
        [ChatToolCall(tool="search_papers", arguments={"query": "video grounding", "limit": 1})],
    )

    assert state.observations[0].status == "completed"
    assert state.references[0]["remote_ingest_token"]
    assert "Video Grounding Paper" in state.context_blocks[0]


@pytest.mark.asyncio
async def test_search_library_uses_rag_service(monkeypatch):
    paper = SimpleNamespace(
        id="paper-1",
        title="Local Paper",
        arxiv_id="2601.00002",
        year=2025,
        authors=["B"],
        abstract="local abstract",
        source_url="https://example.com/local",
        metadata_json={"pdf_url": "https://example.com/local.pdf"},
    )

    class FakeRAG:
        def __init__(self, db):
            self.db = db

        async def search_keyword_and_semantic(self, query, top_k):
            return [(paper, 0.92)]

    monkeypatch.setattr(chat_agent_tools, "RAGService", FakeRAG)

    runtime_state = ChatAgentRuntimeState(user_query="search library")
    runtime_state.db = object()
    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        runtime_state,
        [ChatToolCall(tool="search_library", arguments={"query": "local", "limit": 1})],
    )

    assert state.observations[0].status == "completed"
    assert state.references[0]["source"] == "local_library"
    assert state.references[0]["paper_id"] == "paper-1"


@pytest.mark.asyncio
async def test_read_pdf_uses_full_text_chunks(monkeypatch):
    paper = SimpleNamespace(
        id="paper-1",
        title="Readable Paper",
        arxiv_id="2601.00004",
        year=2026,
        abstract="abstract fallback",
        full_text="full text " * 200,
        source_url="https://example.com/readable",
        metadata_json={"pdf_url": "https://example.com/readable.pdf"},
    )

    async def _load_paper(db, user_id, paper_id):
        return paper

    def _retrieve_chunks(full_text, query, top_k=3):
        return [("method evidence chunk", 0.91), ("experiment evidence chunk", 0.76)], "document"

    monkeypatch.setattr(chat_agent_tools, "_load_user_accessible_paper", _load_paper)
    monkeypatch.setattr(chat_agent_tools.PaperChunkService, "retrieve_chunks", _retrieve_chunks)
    runtime_state = ChatAgentRuntimeState(user_query="read method", user=SimpleNamespace(id="user-1"))
    runtime_state.db = object()

    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        runtime_state,
        [ChatToolCall(tool="read_pdf", arguments={"paper_id": "paper-1", "query": "method", "top_k": 2})],
    )

    assert state.observations[0].status == "completed"
    assert state.observations[0].details["evidence_coverage"] == "full_text"
    assert "method evidence chunk" in state.context_blocks[0]
    assert state.references[0]["source"] == "local_pdf"


@pytest.mark.asyncio
async def test_read_pdf_falls_back_to_abstract(monkeypatch):
    paper = SimpleNamespace(
        id="paper-2",
        title="Abstract Paper",
        arxiv_id=None,
        year=2025,
        abstract="only abstract evidence",
        full_text=None,
        source_url=None,
        metadata_json={},
    )

    async def _load_paper(db, user_id, paper_id):
        return paper

    monkeypatch.setattr(chat_agent_tools, "_load_user_accessible_paper", _load_paper)
    runtime_state = ChatAgentRuntimeState(user_query="read abstract", user=SimpleNamespace(id="user-1"))
    runtime_state.db = object()

    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        runtime_state,
        [ChatToolCall(tool="read_pdf", arguments={"paper_id": "paper-2"})],
    )

    assert state.observations[0].status == "completed"
    assert state.observations[0].details["evidence_coverage"] == "abstract_only"
    assert "only abstract evidence" in state.context_blocks[0]


@pytest.mark.asyncio
async def test_read_pdf_rejects_inaccessible_paper(monkeypatch):
    async def _load_paper(db, user_id, paper_id):
        return None

    monkeypatch.setattr(chat_agent_tools, "_load_user_accessible_paper", _load_paper)
    runtime_state = ChatAgentRuntimeState(user_query="read missing", user=SimpleNamespace(id="user-1"))
    runtime_state.db = object()

    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        runtime_state,
        [ChatToolCall(tool="read_pdf", arguments={"paper_id": "missing"})],
    )

    assert state.observations[0].status == "rejected"
    assert "未找到可访问" in state.observations[0].summary


@pytest.mark.asyncio
async def test_extract_docx_tool_returns_bounded_context():
    from docx import Document

    document = Document()
    document.add_heading("Tool Word Notes", level=1)
    document.add_paragraph("Agent can inspect Word content.")
    buffer = BytesIO()
    document.save(buffer)
    payload = base64.b64encode(buffer.getvalue()).decode("ascii")

    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        ChatAgentRuntimeState(user_query="extract word"),
        [ChatToolCall(tool="extract_docx", arguments={"filename": "notes.docx", "content_base64": payload})],
    )

    observation = state.observations[0]
    assert observation.status == "completed"
    assert observation.details["file_type"] == "docx"
    assert "Tool Word Notes" in state.context_blocks[0]
    assert state.references[0]["source"] == "uploaded_office"


@pytest.mark.asyncio
async def test_extract_pptx_tool_returns_slide_context():
    pytest.importorskip("pptx")
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Tool Slides"
    slide.placeholders[1].text = "First evidence line"
    buffer = BytesIO()
    presentation.save(buffer)
    payload = base64.b64encode(buffer.getvalue()).decode("ascii")

    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        ChatAgentRuntimeState(user_query="extract slides"),
        [ChatToolCall(tool="extract_pptx", arguments={"filename": "slides.pptx", "content_base64": payload})],
    )

    observation = state.observations[0]
    assert observation.status == "completed"
    assert observation.details["file_type"] == "pptx"
    assert observation.details["slide_count"] == 1
    assert "Tool Slides" in state.context_blocks[0]


@pytest.mark.asyncio
async def test_run_skill_tool_returns_structured_context(monkeypatch):
    async def fake_run_skill(skill_id, *, task, context="", current_query="", max_output_chars=4000):
        return ResearchSkillRunResult(
            skill=get_research_skill(skill_id),
            task=task,
            output="实验矩阵：baseline + ablation",
            context_used_chars=len(context),
        )

    monkeypatch.setattr(chat_agent_tools, "run_research_skill", fake_run_skill)

    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        ChatAgentRuntimeState(user_query="用 experiment-planner 做实验设计"),
        [ChatToolCall(tool="run_skill", arguments={"skill_id": "experiment-planner", "task": "设计实验"})],
    )

    observation = state.observations[0]
    assert observation.status == "completed"
    assert observation.details["id"] == "experiment-planner"
    assert observation.references[0]["source"] == "research_skill"
    assert "实验矩阵" in state.context_blocks[0]


@pytest.mark.asyncio
async def test_run_skill_tool_rejects_unknown_skill():
    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        ChatAgentRuntimeState(user_query="run missing skill"),
        [ChatToolCall(tool="run_skill", arguments={"skill_id": "missing-skill", "task": "x"})],
    )

    observation = state.observations[0]
    assert observation.status == "rejected"
    assert "paper-scout" in observation.details["available_skill_ids"]


@pytest.mark.asyncio
async def test_confirmed_import_executes_ingest_and_save(monkeypatch):
    paper = SimpleNamespace(
        id="paper-1",
        title="Imported Paper",
        arxiv_id="2601.00003",
        year=2026,
        source_url="https://arxiv.org/abs/2601.00003",
        metadata_json={"pdf_url": "https://arxiv.org/pdf/2601.00003"},
    )
    calls = {"ingest": 0, "save": 0}

    class FakeIngestion:
        def __init__(self, db):
            self.db = db

        async def ingest_remote(self, source, remote_id, auto_download=False, imported_by_user=None):
            calls["ingest"] += 1
            return paper, True

    class FakeEnhance:
        def __init__(self, db):
            self.db = db

        async def save_paper(self, user_id, paper_id):
            calls["save"] += 1
            return SimpleNamespace(user_id=user_id, paper_id=paper_id)

    monkeypatch.setattr(chat_agent_tools, "PaperIngestionService", FakeIngestion)
    monkeypatch.setattr(chat_agent_tools, "PaperEnhanceService", FakeEnhance)

    args = ImportPaperArgs(source="arxiv", remote_id="2601.00003", auto_download=False).model_dump()
    token = chat_tool_confirmation_token("import_paper", args)
    runtime_state = ChatAgentRuntimeState(user_query="confirm import", user=SimpleNamespace(id="user-1"))
    runtime_state.db = object()
    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        runtime_state,
        [ChatToolCall(tool="import_paper", arguments=args, confirmation_token=token)],
        allow_side_effects=True,
    )

    assert state.observations[0].status == "completed"
    assert calls == {"ingest": 1, "save": 1}
    assert state.references[0]["paper_id"] == "paper-1"


@pytest.mark.asyncio
async def test_confirmed_add_to_folder_executes_membership_mutation():
    folder_id = "11111111-1111-1111-1111-111111111111"
    paper_id = "22222222-2222-2222-2222-222222222222"
    folder = SimpleNamespace(id=UUID(folder_id), name="Grounding")
    paper = SimpleNamespace(
        id=UUID(paper_id),
        title="Folder Paper",
        arxiv_id=None,
        year=2026,
        source_url=None,
        metadata_json={},
    )
    db = _FakeToolDb([[folder], [paper], [], []])
    args = AddToFolderArgs(folder_id=folder_id, paper_ids=[paper_id]).model_dump()
    token = chat_tool_confirmation_token("add_to_folder", args)
    runtime_state = ChatAgentRuntimeState(user_query="confirm folder", user=SimpleNamespace(id="user-1"))
    runtime_state.db = db

    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        runtime_state,
        [ChatToolCall(tool="add_to_folder", arguments=args, confirmation_token=token)],
        allow_side_effects=True,
    )

    assert state.observations[0].status == "completed"
    assert state.observations[0].details["added"] == 1
    assert any(isinstance(item, UserPaper) for item in db.added)
    assert any(isinstance(item, PaperFolderItem) for item in db.added)
    assert db.commit_count == 1


@pytest.mark.asyncio
async def test_confirmed_create_research_project_executes_project_creation():
    paper_id = "22222222-2222-2222-2222-222222222222"
    db = _FakeToolDb([[UUID(paper_id)]])
    args = CreateResearchProjectArgs(
        name="Video Grounding",
        description="Project seed",
        keywords=["video", "grounding"],
        paper_ids=[paper_id],
    ).model_dump()
    token = chat_tool_confirmation_token("create_research_project", args)
    runtime_state = ChatAgentRuntimeState(user_query="confirm project", user=SimpleNamespace(id="user-1"), session_id="session-1")
    runtime_state.db = db

    state = await ChatAgentToolRuntime(default_chat_tool_registry()).run(
        runtime_state,
        [ChatToolCall(tool="create_research_project", arguments=args, confirmation_token=token)],
        allow_side_effects=True,
    )

    assert state.observations[0].status == "completed"
    assert any(isinstance(item, ResearchProject) for item in db.added)
    assert db.added[0].name == "Video Grounding"
    assert db.added[0].paper_ids == [paper_id]
    assert db.commit_count == 1
    assert state.references[0]["source"] == "research_project"


def test_deterministic_import_plan_for_explicit_arxiv_id_waits_confirmation():
    calls = deterministic_chat_tool_plan("请导入 arXiv:2601.00003 到论文库")

    assert calls[0].tool == "import_paper"
    assert calls[0].arguments["source"] == "arxiv"
    assert calls[0].arguments["remote_id"] == "2601.00003"


def test_deterministic_library_plan_prefers_local_library_signal():
    calls = deterministic_chat_tool_plan("请在我的论文库里检索 video grounding")

    assert calls[0].tool == "search_library"


def test_deterministic_plan_routes_local_pdf_reading():
    calls = deterministic_chat_tool_plan("请阅读论文 22222222-2222-2222-2222-222222222222 的方法部分")

    assert calls[0].tool == "read_pdf"
    assert calls[0].arguments["paper_id"] == "22222222-2222-2222-2222-222222222222"


def test_deterministic_plan_routes_folder_and_project_actions():
    folder_calls = deterministic_chat_tool_plan(
        "请把论文 22222222-2222-2222-2222-222222222222 加入分类 11111111-1111-1111-1111-111111111111"
    )
    project_calls = deterministic_chat_tool_plan("创建研究方向 Video Grounding，使用论文 22222222-2222-2222-2222-222222222222")

    assert folder_calls[0].tool == "add_to_folder"
    assert folder_calls[0].arguments["folder_id"] == "11111111-1111-1111-1111-111111111111"
    assert folder_calls[0].arguments["paper_ids"] == ["22222222-2222-2222-2222-222222222222"]
    assert project_calls[0].tool == "create_research_project"
    assert project_calls[0].arguments["paper_ids"] == ["22222222-2222-2222-2222-222222222222"]


def test_deterministic_plan_routes_explicit_skill_prompt():
    calls = deterministic_chat_tool_plan("请用 experiment-planner skill 帮我设计 video grounding 实验")

    assert calls[0].tool == "run_skill"
    assert calls[0].arguments["skill_id"] == "experiment-planner"
