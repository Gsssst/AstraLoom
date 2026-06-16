from io import BytesIO

import pytest

from app.services.office_extraction import (
    UnsupportedOfficeFormatError,
    extract_docx_text,
    extract_office_document,
    extract_pptx_text,
)


def _docx_bytes() -> bytes:
    from docx import Document

    document = Document()
    document.add_heading("Video Grounding Notes", level=1)
    document.add_paragraph("Main paragraph about temporal grounding.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Dataset"
    table.cell(0, 1).text = "Task"
    table.cell(1, 0).text = "Charades-STA"
    table.cell(1, 1).text = "Moment retrieval"
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes() -> bytes:
    pytest.importorskip("pptx")
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Video Grounding"
    slide.placeholders[1].text = "Temporal localization\nMultimodal retrieval"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_docx_extraction_includes_heading_paragraph_and_table():
    result = extract_docx_text(_docx_bytes(), "notes.docx")

    assert result.file_type == "docx"
    assert "[Heading 1] Video Grounding Notes" in result.text
    assert "Main paragraph about temporal grounding." in result.text
    assert "[Table 1]" in result.text
    assert "Charades-STA | Moment retrieval" in result.text
    assert result.metadata["heading_count"] == 1
    assert result.metadata["table_count"] == 1
    assert result.metadata["text_length"] == len(result.text)


def test_pptx_extraction_groups_title_and_text_by_slide():
    result = extract_pptx_text(_pptx_bytes(), "slides.pptx")

    assert result.file_type == "pptx"
    assert "[Slide 1] Video Grounding" in result.text
    assert "Temporal localization" in result.text
    assert "Multimodal retrieval" in result.text
    assert result.metadata["slide_count"] == 1
    assert result.metadata["text_shape_count"] >= 1


def test_legacy_office_formats_return_clear_guidance():
    with pytest.raises(UnsupportedOfficeFormatError) as doc_error:
        extract_office_document(b"legacy", "notes.doc", "application/msword")
    with pytest.raises(UnsupportedOfficeFormatError) as ppt_error:
        extract_office_document(b"legacy", "slides.ppt", "application/vnd.ms-powerpoint")

    assert ".docx" in str(doc_error.value)
    assert ".pptx" in str(ppt_error.value)
