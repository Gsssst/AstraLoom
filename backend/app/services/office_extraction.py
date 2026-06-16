"""Bounded text extraction for Office attachments used by chat tools."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
import re
from typing import Any, Literal


MAX_OFFICE_EXTRACTED_CHARS = 50000
DOCX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
}
PPTX_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
LEGACY_OFFICE_MIME_TYPES = {
    "application/msword",
    "application/vnd.ms-powerpoint",
}


class OfficeExtractionError(ValueError):
    """Raised when an Office file looks supported but cannot be parsed."""


class UnsupportedOfficeFormatError(OfficeExtractionError):
    """Raised when the file is a legacy or unsupported Office format."""


@dataclass(frozen=True)
class OfficeExtractionResult:
    file_type: Literal["docx", "pptx"]
    text: str
    metadata: dict[str, Any]

    @property
    def text_length(self) -> int:
        return len(self.text)


def _filename_suffix(filename: str) -> str:
    return Path(filename or "").suffix.lower()


def is_office_document_candidate(filename: str, content_type: str = "") -> bool:
    suffix = _filename_suffix(filename)
    mime = (content_type or "").split(";")[0].strip().lower()
    return suffix in {".docx", ".doc", ".pptx", ".ppt"} or mime in DOCX_MIME_TYPES | PPTX_MIME_TYPES | LEGACY_OFFICE_MIME_TYPES


def _normalize_text(value: str) -> str:
    return re.sub(r"[ \t\r\f\v]+", " ", (value or "").replace("\xa0", " ")).strip()


def _bounded_text(lines: list[str], max_chars: int) -> tuple[str, bool, int]:
    raw = "\n".join(line for line in lines if line is not None).strip()
    raw_length = len(raw)
    if raw_length <= max_chars:
        return raw, False, raw_length
    suffix = "\n[内容已截断]"
    keep = max(0, max_chars - len(suffix))
    return f"{raw[:keep].rstrip()}{suffix}", True, raw_length


def _docx_heading_level(style_name: str) -> str:
    match = re.search(r"(\d+)", style_name or "")
    return match.group(1) if match else "?"


def _iter_docx_blocks(document: Any):
    from docx.document import Document as DocumentObject
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    if not isinstance(document, DocumentObject):
        return
    parent_elm = document.element.body
    for child in parent_elm.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, document)
        elif child.tag == qn("w:tbl"):
            yield Table(child, document)


def extract_docx_text(
    file_bytes: bytes,
    filename: str = "document.docx",
    *,
    max_chars: int = MAX_OFFICE_EXTRACTED_CHARS,
) -> OfficeExtractionResult:
    try:
        from docx import Document
    except Exception as exc:  # pragma: no cover - dependency is declared.
        raise OfficeExtractionError("Word 解析依赖未安装，请安装 python-docx。") from exc

    try:
        document = Document(BytesIO(file_bytes))
    except Exception as exc:
        raise OfficeExtractionError("Word 文件解析失败，请确认这是有效的 .docx 文件。") from exc

    paragraphs = 0
    headings = 0
    tables = 0
    table_rows = 0
    lines = [f"[Word 文档: {filename}]"]

    for block in _iter_docx_blocks(document):
        block_type = block.__class__.__name__
        if block_type == "Paragraph":
            text = _normalize_text(block.text)
            if not text:
                continue
            paragraphs += 1
            style_name = getattr(getattr(block, "style", None), "name", "") or ""
            if style_name.lower().startswith("heading"):
                headings += 1
                lines.append(f"\n[Heading {_docx_heading_level(style_name)}] {text}")
            else:
                lines.append(text)
        elif block_type == "Table":
            tables += 1
            lines.append(f"\n[Table {tables}]")
            for row in block.rows:
                cells = [_normalize_text(cell.text) for cell in row.cells]
                if any(cells):
                    table_rows += 1
                    lines.append(" | ".join(cells))

    if len(lines) == 1:
        lines.append("未从 Word 文档中提取到可用文本。")

    text, truncated, raw_length = _bounded_text(lines, max_chars)
    return OfficeExtractionResult(
        file_type="docx",
        text=text,
        metadata={
            "filename": filename,
            "file_type": "docx",
            "paragraph_count": paragraphs,
            "heading_count": headings,
            "table_count": tables,
            "table_row_count": table_rows,
            "text_length": len(text),
            "raw_text_length": raw_length,
            "truncated": truncated,
        },
    )


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return _normalize_text(getattr(shape, "text", "") or "")


def _slide_title(slide: Any) -> str:
    for shape in getattr(slide, "shapes", []):
        if not getattr(shape, "is_placeholder", False):
            continue
        text = _shape_text(shape)
        if text:
            return text
    for shape in getattr(slide, "shapes", []):
        text = _shape_text(shape)
        if text:
            return text.splitlines()[0].strip()
    return ""


def extract_pptx_text(
    file_bytes: bytes,
    filename: str = "slides.pptx",
    *,
    max_chars: int = MAX_OFFICE_EXTRACTED_CHARS,
) -> OfficeExtractionResult:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover - dependency may be missing before install.
        raise OfficeExtractionError("PowerPoint 解析依赖未安装，请安装 python-pptx。") from exc

    try:
        presentation = Presentation(BytesIO(file_bytes))
    except Exception as exc:
        raise OfficeExtractionError("PowerPoint 文件解析失败，请确认这是有效的 .pptx 文件。") from exc

    slide_count = len(presentation.slides)
    text_shape_count = 0
    lines = [f"[PowerPoint 演示文稿: {filename}]"]

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title = _slide_title(slide) or f"Slide {slide_index}"
        lines.append(f"\n[Slide {slide_index}] {title}")
        seen_title = False
        slide_text_items = 0
        for shape in slide.shapes:
            text = _shape_text(shape)
            if not text:
                continue
            if text == title and not seen_title:
                seen_title = True
                continue
            text_shape_count += 1
            slide_text_items += 1
            lines.append(f"- {text}")
        if slide_text_items == 0:
            lines.append("- 未提取到正文文本。")

    if slide_count == 0:
        lines.append("未从 PowerPoint 文档中提取到幻灯片文本。")

    text, truncated, raw_length = _bounded_text(lines, max_chars)
    return OfficeExtractionResult(
        file_type="pptx",
        text=text,
        metadata={
            "filename": filename,
            "file_type": "pptx",
            "slide_count": slide_count,
            "text_shape_count": text_shape_count,
            "text_length": len(text),
            "raw_text_length": raw_length,
            "truncated": truncated,
        },
    )


def extract_office_document(
    file_bytes: bytes,
    filename: str,
    content_type: str = "",
    *,
    max_chars: int = MAX_OFFICE_EXTRACTED_CHARS,
) -> OfficeExtractionResult:
    suffix = _filename_suffix(filename)
    mime = (content_type or "").split(";")[0].strip().lower()

    if suffix in {".doc", ".ppt"} or mime in LEGACY_OFFICE_MIME_TYPES:
        raise UnsupportedOfficeFormatError("暂不支持旧版 Office 文件，请另存为 .docx 或 .pptx 后再上传。")
    if suffix == ".docx" or mime in DOCX_MIME_TYPES:
        return extract_docx_text(file_bytes, filename, max_chars=max_chars)
    if suffix == ".pptx" or mime in PPTX_MIME_TYPES:
        return extract_pptx_text(file_bytes, filename, max_chars=max_chars)
    raise UnsupportedOfficeFormatError("不支持的 Office 文件类型，请上传 .docx 或 .pptx 文件。")
